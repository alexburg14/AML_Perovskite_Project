"""Build the project midpoint presentation (presentation/midpoint_presentation.pptx).

Figures are produced by the sibling steps in the chat / from data; re-run after
regenerating figures to refresh the deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
GRAY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width


def slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, text):
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.95)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.25), Inches(12.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def bullets(s, items, left=0.6, top=1.5, width=12.1, height=5.5, size=18):
    tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(6)
        pre = "• " if lvl == 0 else "    – "
        r = p.add_run(); r.text = pre + txt
        r.font.size = Pt(size - 3 * lvl); r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    return tf


def img(s, path, left, top, width):
    try:
        s.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    except Exception as e:
        print("img fail", path, e)


# 1 TITLE
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.2), W, Inches(1.9))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.6)).text_frame
t.word_wrap = True
r = t.paragraphs[0].add_run()
r.text = "Machine-Learning Prediction of Stability and\nBand Gap in ABX₃ Perovskites"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.4)).text_frame
sub.word_wrap = True
r = sub.paragraphs[0].add_run(); r.text = "Project Midpoint Presentation"
r.font.size = Pt(20); r.font.color.rgb = ACCENT
p = sub.add_paragraph(); r = p.add_run(); r.text = "Alexander Burg   ·   29 June 2026"
r.font.size = Pt(16); r.font.color.rgb = GRAY
p = sub.add_paragraph(); r = p.add_run(); r.text = "Atomistic Modeling — M.Sc. project"
r.font.size = Pt(13); r.font.color.rgb = GRAY

# 2 INTRO / RESEARCH QUESTION
s = slide(); title_bar(s, "Introduction & Research Question")
bullets(s, [
    ("ABX₃ perovskites: key materials for solar cells, LEDs, catalysis — but their properties usually come from expensive DFT calculations.", 0),
    ("Goal: predict the two decisive properties directly from composition/structure, bypassing DFT:", 0),
    ("thermodynamic stability — energy above the convex hull (synthesizability)", 1),
    ("electronic structure — band gap (metal vs. semiconductor)", 1),
    ("Constraint: favour synthesizable chemistries — exclude toxic / rare-earth / radioactive elements.", 0),
    ("Research question:", 0),
    ("Can structural ML descriptors predict perovskite stability and band gap well enough to screen novel candidates?", 1),
], width=7.4)
img(s, "analysis/figures/fig10_CsSnI3_polymorphs.png", 8.1, 1.7, 4.9)
cap = s.shapes.add_textbox(Inches(8.1), Inches(5.0), Inches(4.9), Inches(1.4)).text_frame
cap.word_wrap = True
r = cap.paragraphs[0].add_run()
r.text = "One composition (CsSnI₃) → several polymorphs with different band gaps & stability: structural, not just compositional, descriptors are needed."
r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GRAY

# 3 DATA
s = slide(); title_bar(s, "Data: Dataset & Descriptors")
bullets(s, [
    ("Source: OQMD-derived ABX₃ perovskite dataset (Chenebuah & Chenebuah 2023) — 16,323 compounds.", 0),
    ("Relaxed atomic structures fetched from the OQMD API: 16,316 / 16,323 (99.96%).", 0),
    ("Targets per structure: band gap E_g, energy above hull (stability), formation energy.", 0),
    ("Reduced set (remove rare-earth / radioactive): ~8,000 structures; matrix descriptors also cap ≤30 atoms ⇒ 7,875.", 0),
    ("Structural descriptors (DScribe):", 0),
    ("SOAP — 48,870 dimensions (local atomic environments)", 1),
    ("Coulomb matrix — 900 dimensions", 1),
    ("Ewald sum matrix — 900 dimensions (periodicity-aware)", 1),
])

# 4 DATA ANALYSIS - distributions
s = slide(); title_bar(s, "Data Analysis: Target Distributions")
bullets(s, [
    ("Band gap is bimodal — a large spike at 0 eV (metals, ~69%) and the rest above 0.2 eV; almost nothing in between.", 0),
    ("Energy above hull is strongly right-skewed (skew = 2.1): median 0.26, long tail to 4.7 eV/atom.", 0),
    ("Further EDA: element frequencies (84 elements), A/B/X site occupancy, atoms-per-cell (→ n_atoms_max = 30 cutoff).", 0),
], top=1.4, height=1.9, size=16)
img(s, "presentation/figures/fig_targets.png", 1.4, 3.5, 10.5)

# 5 LABELS / SENSITIVITY
s = slide(); title_bar(s, "Data Analysis: Label Definition & Sensitivity")
bullets(s, [
    ("Continuous targets must be thresholded into classes — the cutoff is a modelling choice.", 0),
    ("is_metal = band_gap ≤ t : robust (gaps bimodal) — ~69% metal regardless of t.", 0),
    ("is_stable = E_hull ≤ t : highly sensitive.", 0),
    ("exact-zero ⇒ only 3.3% positive (degenerate task)", 1),
    ("0.025 eV/atom (synthesizability cutoff) ⇒ 13.8%", 1),
    ("⇒ at a 97% majority, plain accuracy is meaningless — use balanced accuracy / PR-AUC.", 0),
], width=7.2, size=16)
img(s, "presentation/figures/fig_label_sensitivity.png", 8.0, 1.8, 5.0)

# 6 ML SETUP
s = slide(); title_bar(s, "Preliminary Tests: Setup")
bullets(s, [
    ("Two tasks: is_metal classification  and  energy-above-hull regression.", 0),
    ("Three descriptors compared: SOAP, Coulomb matrix, Ewald sum matrix.", 0),
    ("Model progression: Decision Tree → Random Forest → linear SVM → HistGradientBoosting.", 0),
    ("Class imbalance handled with class_weight = balanced.", 0),
    ("Evaluation: cross-validation (not a single split); balanced accuracy + F1 for classification, R² + MAE for regression.", 0),
    ("All wired through a small shared package (config / data / descriptors / evaluation / labels).", 0),
])

# 7 RESULTS
s = slide(); title_bar(s, "Preliminary Tests: Results")
bullets(s, [
    ("Regression — best: SOAP + HistGradientBoosting → R² = 0.90, MAE = 0.11 eV/atom.", 0),
    ("Classification — best: SOAP + Random Forest → balanced accuracy = 0.86 (Coulomb only 0.70).", 0),
    ("HistGB is more accurate for regression but 15–35× slower than RF; descriptor choice (SOAP) matters more than model choice for is_metal.", 0),
], top=1.4, height=1.9, size=16)
img(s, "presentation/figures/fig_model_comparison.png", 1.4, 3.5, 10.5)

# 8 CORRELATION WITH DISTRIBUTIONS
s = slide(); title_bar(s, "Does Performance Track the Data Distributions?  Yes.")
bullets(s, [
    ("is_metal: clean bimodal band gap ⇒ well-separated classes ⇒ high, threshold-robust accuracy (0.86).", 0),
    ("is_stable: extreme imbalance (3% positive at exact-zero) ⇒ near-degenerate, low balanced accuracy ⇒ recovered by the 0.025 threshold.", 0),
    ("Hull energy: right-skewed with outliers ⇒ RMSE ≫ MAE; max error ≈ 4.5 eV/atom on rare ultra-unstable phases — the tail is the hard part.", 0),
    ("Dimensionality: SOAP's hull signal is diffuse across 48k dims — RF with max_features=sqrt collapses R² 0.90 → 0.46; the is_metal signal is concentrated (sqrt is fine).", 0),
], size=17)

# 9 LEAKAGE
s = slide(); title_bar(s, "Key Caveat: Honest Evaluation")
bullets(s, [
    ("A random (shuffled) train/test split inflates scores here:", 0),
    ("chemical-family leakage — data is clustered by chemistry", 1),
    ("polymorph leakage — same-composition structures split across train/test", 1),
    ("Holding out grouped/ordered blocks instead:", 0),
    ("hull-energy R²: 0.80 (shuffled) → 0.51 (grouped)", 1),
    ("⇒ realistic screening estimates need composition-grouped splits.", 0),
], width=7.2, size=16)
img(s, "presentation/figures/fig_leakage.png", 8.0, 1.8, 5.0)

# 10 NEXT STEPS
s = slide(); title_bar(s, "Next Steps")
bullets(s, [
    ("Honest evaluation: GroupKFold by composition / scaffold splits to remove leakage.", 0),
    ("Baselines: Dummy + composition-only features — quantify what structure actually adds.", 0),
    ("Dimensionality reduction: PCA(SOAP) → fast and accurate; makes gradient boosting practical at scale.", 0),
    ("Models: HistGradientBoosting (regression) with tuning; stability classifier at 0.025 eV/atom reported with PR-AUC.", 0),
    ("Integrate a toxicity / rare-earth penalty into the screening objective.", 0),
    ("Stretch goal: graph neural network operating directly on the crystal structure.", 0),
])

# 11 REFERENCES
s = slide(); title_bar(s, "References")
refs = [
    "Dataset — E. T. Chenebuah & D. T. Chenebuah, “An inorganic ABX₃ perovskite materials dataset for ML,” arXiv:2312.11335 (2023).  https://doi.org/10.48550/arXiv.2312.11335  ·  github.com/chenebuah/ML_abx3_dataset",
    "OQMD — J. E. Saal et al., JOM 65, 1501 (2013); S. Kirklin et al., npj Comput. Mater. 1, 15010 (2015).  https://oqmd.org",
    "DScribe (descriptors) — L. Himanen et al., Comput. Phys. Commun. 247, 106949 (2020).  https://doi.org/10.1016/j.cpc.2019.106949",
    "SOAP — A. P. Bartók, R. Kondor, G. Csányi, Phys. Rev. B 87, 184115 (2013).",
    "ASE — A. H. Larsen et al., J. Phys.: Condens. Matter 29, 273002 (2017).",
    "scikit-learn — F. Pedregosa et al., J. Mach. Learn. Res. 12, 2825 (2011).",
]
bullets(s, [(x, 0) for x in refs], size=14)

prs.save("presentation/midpoint_presentation.pptx")
print("saved presentation/midpoint_presentation.pptx with", len(prs.slides._sldIdLst), "slides")
